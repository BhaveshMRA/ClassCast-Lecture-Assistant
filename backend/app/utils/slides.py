"""
Slide text extraction helpers.

Supported formats:
  - PPTX  (.pptx)  — extracts text shape-by-shape per slide
  - PDF   (.pdf)   — extracts text page-by-page

Each slide / page is returned as a single string. Slides with no
extractable text are silently skipped.
"""

import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)


def extract_slides_text(file_bytes: bytes, filename: str) -> list[str]:
    """
    Extract text per slide/page from a PPTX or PDF upload.
    Returns a list of strings — one entry per slide/page.
    Raises ValueError for unsupported file types.
    """
    name = filename.lower()
    if name.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    elif name.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file type: '{filename}'. "
            "Please upload a .pptx or .pdf file."
        )


def _extract_pptx(file_bytes: bytes) -> list[str]:
    """Extract text from each slide of a PPTX file."""
    from pptx import Presentation  # type: ignore

    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            # Text frames (titles, body, text boxes)
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        parts.append(line.strip())
            # Tables
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())

        text = "\n".join(parts).strip()
        if text:
            slides.append(text)
            logger.debug(f"slide {i}: {len(text)} chars")
        else:
            logger.debug(f"slide {i}: no extractable text, skipping")

    logger.info(f"pptx: extracted {len(slides)} non-empty slides")
    return slides


def _extract_pdf(file_bytes: bytes) -> list[str]:
    """Extract text from each page of a PDF file."""
    import pdfplumber  # type: ignore

    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                pages.append(text)
                logger.debug(f"pdf page {i}: {len(text)} chars")
            else:
                logger.debug(f"pdf page {i}: no text, skipping")

    logger.info(f"pdf: extracted {len(pages)} non-empty pages")
    return pages
